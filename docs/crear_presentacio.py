#!/usr/bin/env python3
"""Genera Presentacio_PersonalBet.pptx"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "USB_Projecte_PersonalBet" / "05_Presentacio" / "Presentacio_PersonalBet.pptx"
OUT_ALT = Path(__file__).resolve().parent / "Presentacio_PersonalBet.pptx"

BLUE = RGBColor(15, 52, 96)
ORANGE = RGBColor(180, 83, 9)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(30, 41, 59)
GRAY = RGBColor(100, 116, 139)


def set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_title_slide(prs: Presentation, title: str, *lines: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BLUE)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(8.8), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    y = 3.4
    for line in lines:
        b = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(8.8), Inches(0.5))
        t = b.text_frame.paragraphs[0]
        t.text = line
        t.font.size = Pt(18)
        t.font.color.rgb = RGBColor(220, 230, 245)
        t.alignment = PP_ALIGN.CENTER
        y += 0.45


def add_content_slide(prs: Presentation, title: str, bullets: list[str], accent: bool = False) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(248, 250, 252))

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE if accent else BLUE
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.8))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = WHITE

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22 if not bullet.startswith("  ") else 18)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)


def add_arch_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(248, 250, 252))
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.8))
    tb.text_frame.paragraphs[0].text = "Arquitectura"
    tb.text_frame.paragraphs[0].font.size = Pt(28)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = WHITE

    code = """MainActivity
 ├── Apuestas / Añadir apuesta
 ├── Estadísticas (beneficio por fechas)
 ├── Resumen anual (cuentas)
 └── Configuración (CSV)
      ↓
 Room (apuestas) + SharedPreferences"""
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(4.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = code
    p.font.name = "Consolas"
    p.font.size = Pt(20)
    p.font.color.rgb = DARK


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "PersonalBet",
        "Control personal de apuestas deportivas",
        "Alberto Martí Armiñana · DAM",
        "IES Dr. Lluís Simarro · 2025–2026",
        "Tutor: Enric Climent Martí",
    )

    add_content_slide(
        prs,
        "Índice",
        [
            "1. Problema",
            "2. Soluciones anteriores",
            "3. Nuestra solución",
            "4. Beneficio anual y Hacienda",
            "5. Demo de la aplicación",
            "6. Tecnología y arquitectura",
            "7. Resultados y conclusiones",
        ],
    )

    add_content_slide(
        prs,
        "Problema",
        [
            "Usuarios con varias casas de apuestas",
            "Datos dispersos (Excel, notas)",
            "No conocen el ROI ni el beneficio real",
            "Declaración fiscal: difícil calcular el beneficio neto del año",
            "Necesidad: herramienta privada, simple y móvil",
        ],
    )

    add_content_slide(
        prs,
        "Soluciones anteriores",
        [
            "Excel → flexible pero manual y propenso a errores",
            "Apps de casas → solo una operadora",
            "Apps de finanzas → no modelan cuota ni ROI",
            "Vacío: app específica, offline y unificada",
        ],
    )

    add_content_slide(
        prs,
        "Nuestra solución — PersonalBet",
        [
            "App Android nativa en Kotlin",
            "Registro de apuestas + estadísticas + cuentas",
            "Exportación e importación CSV",
            "100 % offline — datos en el dispositivo",
            "Beneficio neto por período (filtro de fechas)",
        ],
    )

    add_content_slide(
        prs,
        "Beneficio anual y Hacienda",
        [
            "Las ganancias en apuestas pueden tributar (AEAT)",
            "Se necesita el beneficio neto del año natural",
            "Solo apuestas ganadas o perdidas (cerradas)",
            "Estadísticas → filtro rango: 01/01 – 31/12",
            "CSV como registro de respaldo",
            "Cifra orientativa — no sustituye asesor fiscal",
        ],
        accent=True,
    )

    add_content_slide(
        prs,
        "Demo (en vivo)",
        [
            "Lista de apuestas y filtros",
            "Añadir apuesta y verificar resultado",
            "Estadísticas → filtrar año 2025 → beneficio neto",
            "Resumen anual: cuentas por casa",
            "Configuración → exportar CSV",
        ],
    )

    add_content_slide(
        prs,
        "Tecnología",
        [
            "Kotlin 2.0 · Room · Coroutines",
            "Material Design · View Binding",
            "Single Activity + 6 Fragments",
            "minSdk 24 · targetSdk 36",
            "Android Studio · Gradle",
        ],
    )

    add_arch_slide(prs)

    add_content_slide(
        prs,
        "Resultados",
        [
            "Objetivos cumplidos (~95 %)",
            "ROI y beneficio neto automáticos",
            "Multi-casa y multi-tipster",
            "Beneficio por ejercicio para referencia fiscal",
            "Privacidad: sin servidor",
        ],
    )

    add_content_slide(
        prs,
        "Limitaciones y mejoras",
        [
            "Solo Android (no iOS ni web)",
            "Sin sincronización en la nube",
            "Importación CSV simplificada",
            "Futuro: MVVM, tests automáticos, iOS",
        ],
    )

    add_content_slide(
        prs,
        "Conclusiones",
        [
            "Problema real → solución técnica viable",
            "Aprendizaje: Android, persistencia, UX",
            "Agradecimientos a Enric Climent Martí",
            "¿Preguntas del tribunal?",
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    prs.save(OUT_ALT)
    print(f"Creat: {OUT}")
    print(f"Copia: {OUT_ALT}")


if __name__ == "__main__":
    main()
